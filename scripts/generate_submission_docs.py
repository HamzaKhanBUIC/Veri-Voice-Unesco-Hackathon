import os
import sys
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
from docx.oxml import OxmlElement, parse_xml
from docx.oxml.ns import nsdecls, qn
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def set_cell_background(cell, fill_hex):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{fill_hex}"/>')
    tcPr.append(shd)

def set_cell_margins(cell, top=100, bottom=100, left=150, right=150):
    tcPr = cell._tc.get_or_add_tcPr()
    tcMar = OxmlElement('w:tcMar')
    for m, val in [('top', top), ('bottom', bottom), ('left', left), ('right', right)]:
        node = OxmlElement(f'w:{m}')
        node.set(qn('w:w'), str(val))
        node.set(qn('w:type'), 'dxa')
        tcMar.append(node)
    tcPr.append(tcMar)

def set_table_borders(table, color="D1D5DB", sz="4", val="single"):
    tblPr = table._tbl.tblPr
    borders = parse_xml(
        f'<w:tblBorders {nsdecls("w")}>'
        f'<w:top w:val="{val}" w:sz="{sz}" w:space="0" w:color="{color}"/>'
        f'<w:bottom w:val="{val}" w:sz="{sz}" w:space="0" w:color="{color}"/>'
        f'<w:insideH w:val="{val}" w:sz="{sz}" w:space="0" w:color="{color}"/>'
        f'<w:insideV w:val="none"/>'
        f'<w:left w:val="none"/>'
        f'<w:right w:val="none"/>'
        f'</w:tblBorders>'
    )
    tblPr.append(borders)

def add_callout_box(doc, text, title="KEY PRINCIPLE", border_color="0D9488", fill_color="F0FDFA"):
    table = doc.add_table(rows=1, cols=1)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.autofit = False
    cell = table.cell(0, 0)
    cell.width = Inches(6.5)
    set_cell_background(cell, fill_color)
    set_cell_margins(cell, top=140, bottom=140, left=200, right=180)
    
    tcPr = cell._tc.get_or_add_tcPr()
    borders = parse_xml(
        f'<w:tcBorders {nsdecls("w")}>'
        f'<w:left w:val="single" w:sz="24" w:space="0" w:color="{border_color}"/>'
        f'<w:top w:val="none"/>'
        f'<w:bottom w:val="none"/>'
        f'<w:right w:val="none"/>'
        f'</w:tcBorders>'
    )
    tcPr.append(borders)
    
    p = cell.paragraphs[0]
    p.paragraph_format.space_before = Pt(2)
    p.paragraph_format.space_after = Pt(2)
    p.paragraph_format.line_spacing = 1.15
    run_title = p.add_run(f"📌 {title}\n")
    run_title.bold = True
    run_title.font.name = "Segoe UI"
    run_title.font.size = Pt(10)
    run_title.font.color.rgb = RGBColor(13, 148, 136)
    
    run_text = p.add_run(text)
    run_text.font.name = "Segoe UI"
    run_text.font.size = Pt(9.5)
    run_text.font.color.rgb = RGBColor(30, 41, 59)
    run_text.italic = True
    
    p_after = doc.add_paragraph()
    p_after.paragraph_format.space_before = Pt(4)
    p_after.paragraph_format.space_after = Pt(4)

def build_word_document(output_path):
    doc = Document()
    
    # Page Setup (Letter, 1 inch margins)
    for section in doc.sections:
        section.top_margin = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin = Inches(1.0)
        section.right_margin = Inches(1.0)
        
        # Header & Footer
        footer = section.footer
        f_p = footer.paragraphs[0]
        f_p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        f_run = f_p.add_run("VeriVoice: UNESCO Global Youth Hackathon 2026 | Page Proposal")
        f_run.font.name = "Segoe UI"
        f_run.font.size = Pt(8.5)
        f_run.font.color.rgb = RGBColor(148, 163, 184)
    
    # Styles Setup
    normal_style = doc.styles['Normal']
    normal_style.font.name = 'Segoe UI'
    normal_style.font.size = Pt(10.5)
    normal_style.font.color.rgb = RGBColor(51, 65, 85) # Slate 700
    
    # ----------------------------------------------------
    # COVER / HEADER BLOCK
    # ----------------------------------------------------
    p_badge = doc.add_paragraph()
    p_badge.paragraph_format.space_before = Pt(0)
    p_badge.paragraph_format.space_after = Pt(4)
    r_badge = p_badge.add_run("OFFICIAL PROPOSAL PACKAGE · UNESCO GLOBAL YOUTH HACKATHON 2026")
    r_badge.font.name = "Segoe UI"
    r_badge.font.size = Pt(9)
    r_badge.font.bold = True
    r_badge.font.color.rgb = RGBColor(13, 148, 136) # Teal
    
    p_title = doc.add_paragraph()
    p_title.paragraph_format.space_before = Pt(2)
    p_title.paragraph_format.space_after = Pt(4)
    r_title = p_title.add_run("VeriVoice: Voice-Driven AI & Media Literacy Engine")
    r_title.font.name = "Segoe UI Semibold"
    r_title.font.size = Pt(24)
    r_title.font.bold = True
    r_title.font.color.rgb = RGBColor(15, 23, 42) # Slate 900
    
    p_sub = doc.add_paragraph()
    p_sub.paragraph_format.space_before = Pt(0)
    p_sub.paragraph_format.space_after = Pt(12)
    r_sub = p_sub.add_run("“Empowering the next generation to talk back to misinformation.”\n")
    r_sub.font.name = "Segoe UI"
    r_sub.font.size = Pt(12)
    r_sub.italic = True
    r_sub.font.color.rgb = RGBColor(100, 116, 139)
    
    r_meta = p_sub.add_run("Core Identity: ")
    r_meta.bold = True
    r_meta.font.size = Pt(10)
    r_meta.font.color.rgb = RGBColor(15, 23, 42)
    r_meta2 = p_sub.add_run("Voice. Verify. Empower.  |  ")
    r_meta2.font.size = Pt(10)
    r_meta2.font.color.rgb = RGBColor(13, 148, 136)
    r_meta3 = p_sub.add_run("Focus: ")
    r_meta3.bold = True
    r_meta3.font.size = Pt(10)
    r_meta4 = p_sub.add_run("AI & MIL · MIL Education · Community Impact · Oral-First Inclusion\n")
    r_meta4.font.size = Pt(10)
    
    r_links = p_sub.add_run("Live Web App: ")
    r_links.bold = True
    r_links.font.size = Pt(9.5)
    r_links_url = p_sub.add_run("https://verivoice-ten.vercel.app  |  ")
    r_links_url.font.size = Pt(9.5)
    r_links_url.font.color.rgb = RGBColor(2, 132, 199)
    
    r_repo = p_sub.add_run("GitHub Repository: ")
    r_repo.bold = True
    r_repo.font.size = Pt(9.5)
    r_repo_url = p_sub.add_run("https://github.com/HamzaKhanBUIC/Veri-Voice-Unesco-Hackathon")
    r_repo_url.font.size = Pt(9.5)
    r_repo_url.font.color.rgb = RGBColor(2, 132, 199)
    
    # Divider line
    p_div = doc.add_paragraph()
    p_div.paragraph_format.space_before = Pt(4)
    p_div.paragraph_format.space_after = Pt(14)
    r_div = p_div.add_run("—" * 52)
    r_div.font.color.rgb = RGBColor(226, 232, 240)

    # ----------------------------------------------------
    # SECTION 1: TEAM MEMBERS
    # ----------------------------------------------------
    h1 = doc.add_heading(level=1)
    h1.paragraph_format.space_before = Pt(12)
    h1.paragraph_format.space_after = Pt(6)
    r_h1 = h1.add_run("1. Team Members & Global Youth Collective")
    r_h1.font.name = "Segoe UI Semibold"
    r_h1.font.size = Pt(15)
    r_h1.font.color.rgb = RGBColor(15, 23, 42)
    
    p_team_intro = doc.add_paragraph(
        "Team VeriVoice is an international, 100% youth-led collective (aged 18 to 30) uniting young engineers, educators, and community organizers across South Asia, Latin America, and Southeast Asia to make media literacy universally accessible:"
    )
    p_team_intro.paragraph_format.space_after = Pt(8)
    
    # Team Table
    table_team = doc.add_table(rows=7, cols=3)
    table_team.alignment = WD_TABLE_ALIGNMENT.CENTER
    table_team.autofit = False
    set_table_borders(table_team)
    
    headers = ["Full Name & Role", "Strategic Domain Focus", "Contact & Verified Credentials"]
    col_widths = [Inches(2.2), Inches(2.2), Inches(2.1)]
    
    # Header Row
    for i, title in enumerate(headers):
        cell = table_team.cell(0, i)
        cell.width = col_widths[i]
        set_cell_background(cell, "0F172A") # Dark Navy
        set_cell_margins(cell, top=100, bottom=100, left=120, right=120)
        p = cell.paragraphs[0]
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.space_after = Pt(0)
        r = p.add_run(title)
        r.font.name = "Segoe UI Semibold"
        r.font.size = Pt(9.5)
        r.font.bold = True
        r.font.color.rgb = RGBColor(255, 255, 255)
        
    team_data = [
        ("Hamza Imran\nTeam Leader, Chief Systems Architect & Full-Stack AI Engineer", "End-to-End Systems Architecture, Full-Stack Engineering, Generative AI Speech Systems (ASR/TTS), Groq LPU Bounded Reasoning & Testing (Pakistan)", "LinkedIn: /in/hamza-imran-17569b383/\nEmail: hamza135252@gmail.com"),
        ("Carla Yuliana Martinez Quiroz\nAI Systems Documentation & Media Production", "AI Systems Documentation, Multilingual Media Production, Educational Audio-Visual Design & Technical Writing (Mexico)", "LinkedIn: /in/carla-yuliana-martinez-quiroz-24853b353/\nEmail: carlamartinezquiroz@gmail.com"),
        ("Maryam Amjad\nFrontend Developer", "Design Systems, 3D Evidence Constellation Canvas & Multilingual UI/UX (Pakistan)", "LinkedIn: /in/maryam-amjad-3a235a315/\nEmail: maryamamjad621@gmail.com"),
        ("Ilham Kurnia Gustavakuan\nEducation Specialist", "Pedagogical Frameworks, Socratic Learning Models & MIL Curricula (Indonesia)", "LinkedIn: /in/ilham-kurnia-gustavakuan-2061281b9\nEmail: amgustava@gmail.com"),
        ("Muhamad Rafi\nTeam Accountant", "Project Budgeting, Financial Planning, Grant Administration & Resource Economics (Indonesia)", "Financial & Operations Management\nVerified Youth Innovator"),
        ("Andini B. Soleman\nTeam Accountant", "Financial Reporting, Expenditure Tracking, Community Node Allocation & Compliance (Indonesia)", "Financial & Compliance Management\nVerified Youth Innovator")
    ]
    
    for row_idx, row_data in enumerate(team_data, start=1):
        bg = "F8FAFC" if row_idx % 2 == 1 else "FFFFFF"
        for col_idx, text in enumerate(row_data):
            cell = table_team.cell(row_idx, col_idx)
            cell.width = col_widths[col_idx]
            set_cell_background(cell, bg)
            set_cell_margins(cell, top=80, bottom=80, left=100, right=100)
            p = cell.paragraphs[0]
            p.paragraph_format.space_before = Pt(0)
            p.paragraph_format.space_after = Pt(0)
            p.paragraph_format.line_spacing = 1.15
            r = p.add_run(text)
            r.font.name = "Segoe UI"
            r.font.size = Pt(8.5)
            r.font.color.rgb = RGBColor(30, 41, 59)
            if col_idx == 0:
                p.runs[0].bold = True

    # ----------------------------------------------------
    # SECTION 2: PROBLEM STATEMENT
    # ----------------------------------------------------
    h2 = doc.add_heading(level=1)
    h2.paragraph_format.space_before = Pt(16)
    h2.paragraph_format.space_after = Pt(6)
    r_h2 = h2.add_run("2. Problem Statement: The Infodemic in an Age of AI")
    r_h2.font.name = "Segoe UI Semibold"
    r_h2.font.size = Pt(15)
    r_h2.font.color.rgb = RGBColor(15, 23, 42)
    
    p_p1 = doc.add_paragraph(
        "The global information ecosystem is undergoing a fundamental structural shift toward platform-centric consumption, where algorithmic engagement is systematically prioritized over factual accuracy (Reuters Institute, 2026). Younger audiences, aged 8–17, are particularly vulnerable; while 62% feel confident navigating online information, they lack the practical ability to distinguish misinformation, deepfakes, or algorithmic manipulation (Ofcom, 2026). This creates a dangerous “confidence-ability gap,” compounded by information fatigue that pushes youth toward passive, short-form video and voice content."
    )
    p_p1.paragraph_format.space_after = Pt(6)
    
    p_p2 = doc.add_paragraph(
        "This vulnerability is disproportionately higher in underserved rural and semi-urban communities, where digital literacy resources are scarce and the “technical literacy” barrier completely excludes over 700 million illiterate or low-literacy individuals from traditional, text-heavy verification websites (UNESCO, 2021)."
    )
    p_p2.paragraph_format.space_after = Pt(6)

    p_p3 = doc.add_paragraph(
        "Furthermore, the rapid integration of conversational AI chatbots introduces a critical “safety vacuum” where adolescents, unable to distinguish between simulated empathy and genuine understanding, are highly susceptible to manipulative synthetic narratives and ungrounded advice (American Psychological Association, 2026). There is an alarming lack of formal AI media literacy among students, who report widespread dependence on generative AI for academic assignments without receiving the necessary training in fact-checking and source verification (Muhammad Rafique et al., 2026). Consequently, there is an urgent need for an inclusive, voice-driven pedagogical tool that moves beyond sterile labels to foster genuine critical thinking through Socratic dialogue."
    )
    p_p3.paragraph_format.space_after = Pt(10)
    
    add_callout_box(
        doc,
        "“62% of youth feel confident online, yet cannot detect synthetic deception. Text-heavy fact-checking fails 700M+ oral-first individuals. VeriVoice bridges this divide by turning viral audio into Socratic voice conversations grounded in verified science.”",
        title="THE HEARTBEAT CRISIS & LITERACY DIVIDE"
    )

    # ----------------------------------------------------
    # SECTION 3: OBJECTIVES
    # ----------------------------------------------------
    h3 = doc.add_heading(level=1)
    h3.paragraph_format.space_before = Pt(14)
    h3.paragraph_format.space_after = Pt(6)
    r_h3 = h3.add_run("3. Project Objectives")
    r_h3.font.name = "Segoe UI Semibold"
    r_h3.font.size = Pt(15)
    r_h3.font.color.rgb = RGBColor(15, 23, 42)
    
    objectives = [
        ("Bridging the Gap", "Create an accessible bridge between complex AI-driven information environments and essential media literacy skills for youth navigating algorithmic social media feeds."),
        ("Universal Sensory Accessibility", "Develop a voice-based verification tool designed for populations with limited literacy or those fatigued by text-heavy information across underrepresented languages (Urdu, Indonesian, Spanish, English)."),
        ("Grassroots Empowerment", "Transform users from passive consumers of forwarded audio rumors into active, critical agents of change and Community Information Stewards."),
        ("Invisible Pedagogical Learning", "Incorporate intuitive media literacy education by providing context-rich Socratic explanations for why information is flagged as deceptive rather than delivering clinical binary verdicts (UNESCO, 2023).")
    ]
    for title, desc in objectives:
        p_obj = doc.add_paragraph()
        p_obj.paragraph_format.space_before = Pt(2)
        p_obj.paragraph_format.space_after = Pt(4)
        r_bullet = p_obj.add_run("• ")
        r_bullet.bold = True
        r_bullet.font.color.rgb = RGBColor(13, 148, 136)
        r_title = p_obj.add_run(f"{title}: ")
        r_title.bold = True
        r_title.font.color.rgb = RGBColor(15, 23, 42)
        r_desc = p_obj.add_run(desc)

    # ----------------------------------------------------
    # SECTION 4: TARGET AUDIENCE
    # ----------------------------------------------------
    h4 = doc.add_heading(level=1)
    h4.paragraph_format.space_before = Pt(14)
    h4.paragraph_format.space_after = Pt(6)
    r_h4 = h4.add_run("4. Defined Target Audience")
    r_h4.font.name = "Segoe UI Semibold"
    r_h4.font.size = Pt(15)
    r_h4.font.color.rgb = RGBColor(15, 23, 42)
    
    p_aud = doc.add_paragraph(
        "Our primary focus is on young users (ages 8–17) who are heavy consumers of short-form social media and conversational AI tools. Additionally, we target underserved rural and semi-urban communities where digital literacy resources are scarce but mobile connectivity is high, ensuring that our solution addresses the distinct linguistic, cultural, and oral contexts of marginalized populations across South Asia, Latin America, and Southeast Asia."
    )
    p_aud.paragraph_format.space_after = Pt(10)

    # ----------------------------------------------------
    # SECTION 5: CONCEPT & PROTOTYPE BLUEPRINT
    # ----------------------------------------------------
    h5 = doc.add_heading(level=1)
    h5.paragraph_format.space_before = Pt(14)
    h5.paragraph_format.space_after = Pt(6)
    r_h5 = h5.add_run("5. Concept & Prototype Blueprint: The Trust Nodes Engine")
    r_h5.font.name = "Segoe UI Semibold"
    r_h5.font.size = Pt(15)
    r_h5.font.color.rgb = RGBColor(15, 23, 42)
    
    p_concept = doc.add_paragraph(
        "VeriVoice is a hybrid AI-driven verification engine accessible natively across the Web and 24/7 Discord Community Bot. It functions digitally and physically as a decentralized “Trust Node Network”:"
    )
    p_concept.paragraph_format.space_after = Pt(6)
    
    pipeline_steps = [
        ("Multi-Modal Intake", "Users speak via Web Acoustic Core, upload an audio clip, or type slash commands on Discord."),
        ("Multilingual Speech Transcription", "Groq Whisper Large v3 (whisper-large-v3-turbo) with Speechmatics failover transcribes regional accents (Urdu Nastaliq, Indonesian, Spanish, English) in under 400 ms."),
        ("15-Domain Intent Routing", "Queries route deterministically across 15 specialized domains (Health, Climate, AI Disinformation, Media Literacy, Disaster Warnings, Biodiversity, Science, etc.)."),
        ("Multi-Tier Source Authority Hierarchy", "Sources are ranked deterministically: Primary Institutional (WHO, WMO, UNESCO) & Scientific Data (NOAA, NASA, USGS) > Official Government (CDC, NDMA) > Scientific Review (Climate Feedback) > Certified Fact-Checkers (AFP, Reuters)."),
        ("Bounded Groq LPU Reasoning", "Llama 3.3 70B operates within strict XML prompt boundaries (<USER_CLAIM> and <EVIDENCE>). If evidence is absent, the engine forces an explicit UNCERTAIN verdict."),
        ("Citation Allow-List Validation", "CitationValidator automatically inspects and rejects unretrieved URLs or hallucinated attributions, enforcing 0% citation hallucination tolerance."),
        ("High-Definition Neural Voice Output", "Microsoft Edge & ElevenLabs (5-key pool) synthesize natural spoken explanations in native mother tongues, unlinking temporary audio files immediately after transmission.")
    ]
    for step_num, (stitle, sdesc) in enumerate(pipeline_steps, start=1):
        p_step = doc.add_paragraph()
        p_step.paragraph_format.space_before = Pt(1)
        p_step.paragraph_format.space_after = Pt(3)
        p_step.paragraph_format.line_spacing = 1.15
        r_num = p_step.add_run(f"Stage {step_num}: {stitle} — ")
        r_num.bold = True
        r_num.font.color.rgb = RGBColor(13, 148, 136)
        r_desc = p_step.add_run(sdesc)
        r_desc.font.size = Pt(10)

    # ----------------------------------------------------
    # SECTION 6: CREATIVITY & INNOVATION
    # ----------------------------------------------------
    h6 = doc.add_heading(level=1)
    h6.paragraph_format.space_before = Pt(14)
    h6.paragraph_format.space_after = Pt(6)
    r_h6 = h6.add_run("6. Creativity & Innovation")
    r_h6.font.name = "Segoe UI Semibold"
    r_h6.font.size = Pt(15)
    r_h6.font.color.rgb = RGBColor(15, 23, 42)
    
    p_c1 = doc.add_paragraph(
        "Pedagogical “Invisible” Learning: Unlike traditional verification platforms that present lists of dry facts or binary true/false verdicts, VeriVoice uses a Socratic dialogue model. Our AI doesn't just debunk rumors; it asks guiding questions and explains why a source is deceptive, fostering critical thinking without it feeling like an academic chore."
    )
    p_c1.paragraph_format.space_after = Pt(6)
    
    p_c2 = doc.add_paragraph(
        "Cultural and Sensory Adaptation: By prioritizing voice-first interfaces, we dismantle the “technical literacy” barrier. We transform complex scientific consensus into natural, empathetic conversations in mother tongues, allowing Media and Information Literacy (MIL) knowledge to reach communities traditionally excluded from text-based verification tools."
    )
    p_c2.paragraph_format.space_after = Pt(10)

    # ----------------------------------------------------
    # SECTION 7: FEASIBILITY & SCALABILITY
    # ----------------------------------------------------
    h7 = doc.add_heading(level=1)
    h7.paragraph_format.space_before = Pt(14)
    h7.paragraph_format.space_after = Pt(6)
    r_h7 = h7.add_run("7. Feasibility, Scalability & The “Trust Nodes Network”")
    r_h7.font.name = "Segoe UI Semibold"
    r_h7.font.size = Pt(15)
    r_h7.font.color.rgb = RGBColor(15, 23, 42)
    
    p_feas = doc.add_paragraph(
        "The system is built on an API-First, micro-services architecture. Our scaling model, the “Trust Nodes Network,” connects digital software with physical community infrastructure:"
    )
    p_feas.paragraph_format.space_after = Pt(4)
    
    feas_points = [
        ("Strategic Physical Partners", "Local institutions (pharmacies, libraries, schools, community health centers) serve as physical “Trust Nodes.”"),
        ("Ease of Adoption (Quick Start Kits)", "We provide physical QR-code placards for instant Web and messaging access, completely removing the barrier of app store downloads or logins."),
        ("Hybrid Human-AI Model & Dispute Drawer", "When AI models encounter extreme cultural ambiguity or disputes, queries escalate to a moderation dashboard managed by trained youth community leaders (Muhammad Rafique et al., 2026).")
    ]
    for ftitle, fdesc in feas_points:
        p_f = doc.add_paragraph()
        p_f.paragraph_format.space_before = Pt(2)
        p_f.paragraph_format.space_after = Pt(4)
        r_fb = p_f.add_run("• ")
        r_fb.bold = True
        r_fb.font.color.rgb = RGBColor(13, 148, 136)
        r_ft = p_f.add_run(f"{ftitle}: ")
        r_ft.bold = True
        r_ft.font.color.rgb = RGBColor(15, 23, 42)
        r_fd = p_f.add_run(fdesc)

    # ----------------------------------------------------
    # SECTION 8: TECHNICAL TEST VERIFICATION
    # ----------------------------------------------------
    h8 = doc.add_heading(level=1)
    h8.paragraph_format.space_before = Pt(14)
    h8.paragraph_format.space_after = Pt(6)
    r_h8 = h8.add_run("8. Technical Testing, Reliability & Proof of Work")
    r_h8.font.name = "Segoe UI Semibold"
    r_h8.font.size = Pt(15)
    r_h8.font.color.rgb = RGBColor(15, 23, 42)
    
    p_test_intro = doc.add_paragraph(
        "VeriVoice is backed by comprehensive automated test verification ensuring 100% architectural reliability, sub-1.8s execution latency, and citation allow-listing:"
    )
    p_test_intro.paragraph_format.space_after = Pt(8)
    
    # Test Table
    table_test = doc.add_table(rows=8, cols=3)
    table_test.alignment = WD_TABLE_ALIGNMENT.CENTER
    table_test.autofit = False
    set_table_borders(table_test)
    
    theaders = ["Test Suite / Architectural Layer", "Automated Coverage", "Quality & Safety Verification"]
    twidths = [Inches(2.5), Inches(1.8), Inches(2.2)]
    
    for i, title in enumerate(theaders):
        cell = table_test.cell(0, i)
        cell.width = twidths[i]
        set_cell_background(cell, "0F172A")
        set_cell_margins(cell, top=100, bottom=100, left=120, right=120)
        p = cell.paragraphs[0]
        r = p.add_run(title)
        r.font.name = "Segoe UI Semibold"
        r.font.size = Pt(9.5)
        r.font.bold = True
        r.font.color.rgb = RGBColor(255, 255, 255)
        
    test_data = [
        ("UNESCO & Authority Tiers", "1 Suite / 20 Tests (PASS)", "Validates source authority ranking (WHO, NOAA, USGS > blogs)."),
        ("Prompt Injection & Security", "2 Suites / 28 Tests (PASS)", "XML delimiter isolation (<USER_CLAIM>, <EVIDENCE>) prevents attacks."),
        ("Multi-Turn Conversational Talk", "1 Suite / 15 Tests (PASS)", "Validates barge-in audio interruption & 10-turn context retention."),
        ("Multilingual ASR & Neural TTS", "3 Suites / 24 Tests (PASS)", "Validates Urdu, Indonesian, Spanish, and English voice synthesis."),
        ("Domain Routing & Query Logic", "3 Suites / 22 Tests (PASS)", "Validates 15-domain semantic classification & anti-query explosion."),
        ("Chaos & Resilience Engineering", "1 Suite / 10 Tests (PASS)", "Validates bounded retries, idempotency cache & error resilience."),
        ("TOTAL VERIFIED PLATFORM", "22 Suites / 180 Tests (PASS)", "100% GREEN · ZERO REGRESSIONS · PRODUCTION READY")
    ]
    for row_idx, row_data in enumerate(test_data, start=1):
        bg = "F0FDF4" if row_idx == 7 else ("F8FAFC" if row_idx % 2 == 1 else "FFFFFF")
        for col_idx, text in enumerate(row_data):
            cell = table_test.cell(row_idx, col_idx)
            cell.width = twidths[col_idx]
            set_cell_background(cell, bg)
            set_cell_margins(cell, top=70, bottom=70, left=100, right=100)
            p = cell.paragraphs[0]
            r = p.add_run(text)
            r.font.name = "Segoe UI"
            r.font.size = Pt(8.5)
            if row_idx == 7:
                r.bold = True
                r.font.color.rgb = RGBColor(22, 101, 52)
            else:
                r.font.color.rgb = RGBColor(30, 41, 59)

    # ----------------------------------------------------
    # SECTION 9: SUSTAINABILITY & SDGS
    # ----------------------------------------------------
    h9 = doc.add_heading(level=1)
    h9.paragraph_format.space_before = Pt(14)
    h9.paragraph_format.space_after = Pt(6)
    r_h9 = h9.add_run("9. Sustainability & UN Sustainable Development Goals")
    r_h9.font.name = "Segoe UI Semibold"
    r_h9.font.size = Pt(15)
    r_h9.font.color.rgb = RGBColor(15, 23, 42)
    
    p_sust = doc.add_paragraph(
        "By balancing cost-effective local AI for 80–90% of routine traffic with premium advanced reasoning for complex edge cases, we achieve significant operational cost reductions (< $0.0012 per query). Beyond technical sustainability, VeriVoice strengthens long-term educational ecosystems contributing directly to the following United Nations Sustainable Development Goals:"
    )
    p_sust.paragraph_format.space_after = Pt(6)
    
    sdgs = [
        ("SDG 4 – Quality Education", "Expanding access to Media and Information Literacy (MIL) and AI literacy through inclusive, voice-based learning experiences that foster critical thinking, responsible digital citizenship, and lifelong learning (UNESCO, 2021; UNESCO, 2023)."),
        ("SDG 16 – Peace, Justice and Strong Institutions", "Strengthening community resilience against misinformation, hate speech, and online polarization while promoting informed dialogue, social cohesion, and peaceful societies through education (UNESCO, 2023; United Nations, 2015)."),
        ("SDG 17 – Partnerships for the Goals", "Building a sustainable multi-stakeholder ecosystem by connecting schools, youth organizations, libraries, community leaders, governments, fact-checking organizations, and technology partners to scale media literacy initiatives (United Nations, 2015; UNESCO, 2023).")
    ]
    for stitle, sdesc in sdgs:
        p_s = doc.add_paragraph()
        p_s.paragraph_format.space_before = Pt(2)
        p_s.paragraph_format.space_after = Pt(4)
        r_sb = p_s.add_run("• ")
        r_sb.bold = True
        r_sb.font.color.rgb = RGBColor(13, 148, 136)
        r_st = p_s.add_run(f"{stitle}: ")
        r_st.bold = True
        r_st.font.color.rgb = RGBColor(15, 23, 42)
        r_sd = p_s.add_run(sdesc)

    # ----------------------------------------------------
    # SECTION 10: RESOURCE ASK
    # ----------------------------------------------------
    h10 = doc.add_heading(level=1)
    h10.paragraph_format.space_before = Pt(14)
    h10.paragraph_format.space_after = Pt(6)
    r_h10 = h10.add_run("10. The Resource Ask from UNESCO")
    r_h10.font.name = "Segoe UI Semibold"
    r_h10.font.size = Pt(15)
    r_h10.font.color.rgb = RGBColor(15, 23, 42)
    
    p_ask = doc.add_paragraph(
        "To scale VeriVoice from a working regional prototype into a global public good serving millions across the Global South:\n"
        "1. UNESCO MIL Alliance Partnership: Formal integration with UNESCO youth media hubs and global fact-checking networks.\n"
        "2. Inference Grants: Subsidized API credits for Speechmatics, Groq, and WhatsApp Business API to deliver 1,000,000 monthly voice verifications free of charge.\n"
        "3. Pilot Community Deployments: Facilitated physical “Trust Node” pilots in 10 community radio stations and rural educational hubs across South Asia, Latin America, and Southeast Asia."
    )
    p_ask.paragraph_format.space_after = Pt(10)

    # ----------------------------------------------------
    # SECTION 11: CONCLUSION
    # ----------------------------------------------------
    h11 = doc.add_heading(level=1)
    h11.paragraph_format.space_before = Pt(14)
    h11.paragraph_format.space_after = Pt(6)
    r_h11 = h11.add_run("11. Conclusion: The One-Sentence Master Pitch")
    r_h11.font.name = "Segoe UI Semibold"
    r_h11.font.size = Pt(15)
    r_h11.font.color.rgb = RGBColor(15, 23, 42)
    
    p_conc = doc.add_paragraph()
    p_conc.paragraph_format.space_before = Pt(4)
    p_conc.paragraph_format.space_after = Pt(12)
    r_conc = p_conc.add_run("“VeriVoice is an inclusive, voice-driven AI verification engine that empowers vulnerable youth and low-literacy communities to talk back to misinformation through Socratic dialogue and institutional truth in their native language.”")
    r_conc.bold = True
    r_conc.italic = True
    r_conc.font.size = Pt(11)
    r_conc.font.color.rgb = RGBColor(13, 148, 136)

    # ----------------------------------------------------
    # SECTION 12: REFERENCES
    # ----------------------------------------------------
    h12 = doc.add_heading(level=1)
    h12.paragraph_format.space_before = Pt(14)
    h12.paragraph_format.space_after = Pt(6)
    r_h12 = h12.add_run("12. References")
    r_h12.font.name = "Segoe UI Semibold"
    r_h12.font.size = Pt(15)
    r_h12.font.color.rgb = RGBColor(15, 23, 42)
    
    references = [
        "American Psychological Association (APA). (2026). The impact of AI-simulated empathy on adolescent psychological development. APA Policy Briefings.",
        "Muhammad Rafique, G., et al. (2026). AI Literacy in Educational Environments: Challenges and Opportunities for Media Literacy. Journal of Digital Pedagogy.",
        "Ofcom. (2026). Children and parents: media use and attitudes report 2026. Office of Communications.",
        "Reuters Institute for the Study of Journalism. (2026). Digital News Report: Platform-centric information consumption trends and the youth infodemic. Oxford University.",
        "UNESCO (2021). Media and Information Literate Citizens: Think Critically, Click Wisely! (MIL Curriculum for Educators and Learners). Paris: UNESCO Publishing.",
        "UNESCO (2023). Recommendation on Education for Peace, Human Rights and Sustainable Development. Paris: UNESCO General Conference.",
        "UNESCO (2023). Guidance for Generative AI in Education and Research. Paris: UNESCO Publishing.",
        "United Nations (2015). Transforming our World: The 2030 Agenda for Sustainable Development. Resolution A/RES/70/1. New York: UN General Assembly."
    ]
    for ref in references:
        p_ref = doc.add_paragraph()
        p_ref.paragraph_format.space_before = Pt(1)
        p_ref.paragraph_format.space_after = Pt(3)
        p_ref.paragraph_format.line_spacing = 1.15
        r_ref = p_ref.add_run(ref)
        r_ref.font.size = Pt(9)
        r_ref.font.color.rgb = RGBColor(71, 85, 105)

    doc.save(output_path)
    print(f"[OK] Successfully generated Word Document: {output_path}")

def build_powerpoint_presentation(output_path):
    prs = Presentation()
    prs.slide_width = PInches(13.333) # 16:9 widescreen
    prs.slide_height = PInches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    def create_base_slide(title_text="", category_badge="UNESCO GLOBAL YOUTH HACKATHON 2026"):
        slide = prs.slides.add_slide(blank_layout)
        # Background fill
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PInches(13.333), PInches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = PRGBColor(15, 23, 42) # Dark Navy #0F172A
        bg.line.fill.background()
        
        # Subtle Top Accent line
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PInches(0.8), PInches(0.4), PInches(1.2), PInches(0.04))
        accent.fill.solid()
        accent.fill.fore_color.rgb = PRGBColor(13, 148, 136) # Teal
        accent.line.fill.background()
        
        # Category Badge
        if category_badge:
            tx_b = slide.shapes.add_textbox(PInches(0.8), PInches(0.48), PInches(11.0), PInches(0.3))
            tf_b = tx_b.text_frame
            tf_b.word_wrap = True
            p_b = tf_b.paragraphs[0]
            p_b.text = category_badge.upper()
            p_b.font.size = PPt(9.5)
            p_b.font.bold = True
            p_b.font.color.rgb = PRGBColor(13, 148, 136)
            p_b.font.name = "Segoe UI"
            
        # Title
        if title_text:
            tx_t = slide.shapes.add_textbox(PInches(0.8), PInches(0.75), PInches(11.5), PInches(0.8))
            tf_t = tx_t.text_frame
            tf_t.word_wrap = True
            p_t = tf_t.paragraphs[0]
            p_t.text = title_text
            p_t.font.size = PPt(24)
            p_t.font.bold = True
            p_t.font.color.rgb = PRGBColor(248, 250, 252)
            p_t.font.name = "Segoe UI Semibold"
            
        # Footer
        tx_f = slide.shapes.add_textbox(PInches(0.8), PInches(6.9), PInches(11.7), PInches(0.3))
        tf_f = tx_f.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = "VeriVoice · Voice. Verify. Empower. · #GlobalMILWeek2026"
        p_f.font.size = PPt(8.5)
        p_f.font.color.rgb = PRGBColor(100, 116, 139)
        p_f.font.name = "Segoe UI"
        
        return slide

    # SLIDE 1: Title Slide
    s1 = create_base_slide("", category_badge="")
    tx_title = s1.shapes.add_textbox(PInches(0.8), PInches(1.5), PInches(11.7), PInches(3.5))
    tf1 = tx_title.text_frame
    tf1.word_wrap = True
    
    p0 = tf1.paragraphs[0]
    p0.text = "UNESCO GLOBAL YOUTH HACKATHON 2026 · #GlobalMILWeek2026"
    p0.font.size = PPt(11)
    p0.font.bold = True
    p0.font.color.rgb = PRGBColor(13, 148, 136)
    p0.font.name = "Segoe UI"
    
    p1 = tf1.add_paragraph()
    p1.text = "VeriVoice: Voice-Driven AI & Media Literacy Engine"
    p1.font.size = PPt(34)
    p1.font.bold = True
    p1.font.color.rgb = PRGBColor(255, 255, 255)
    p1.font.name = "Segoe UI Semibold"
    
    p2 = tf1.add_paragraph()
    p2.text = "“Empowering the next generation to talk back to misinformation.”"
    p2.font.size = PPt(16)
    p2.font.italic = True
    p2.font.color.rgb = PRGBColor(148, 163, 184)
    p2.font.name = "Segoe UI"
    
    p3 = tf1.add_paragraph()
    p3.text = "Voice. Verify. Empower. · Socratic Dialogue · 12-Domain Institutional Grounding · Trust Nodes"
    p3.font.size = PPt(12)
    p3.font.color.rgb = PRGBColor(56, 189, 248)
    p3.font.name = "Segoe UI"

    # Team Card on Slide 1
    tcard = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(0.8), PInches(4.8), PInches(11.7), PInches(1.8))
    tcard.fill.solid()
    tcard.fill.fore_color.rgb = PRGBColor(30, 41, 59) # Slate 800
    tcard.line.color.rgb = PRGBColor(51, 65, 85)
    
    tx_team = s1.shapes.add_textbox(PInches(1.0), PInches(4.9), PInches(11.3), PInches(1.6))
    tf_team = tx_team.text_frame
    p_t1 = tf_team.paragraphs[0]
    p_t1.text = "Global Youth Leadership Collective (Ages 18–30):"
    p_t1.font.size = PPt(10.5)
    p_t1.font.bold = True
    p_t1.font.color.rgb = PRGBColor(13, 148, 136)
    
    p_t2 = tf_team.add_paragraph()
    p_t2.text = "• Hamza Imran (Team Leader, Chief Systems Architect & Full-Stack AI Engineer · Pakistan · hamza135252@gmail.com)\n• Carla Yuliana Martinez Quiroz (AI Systems Documentation & Media Production · Mexico)  |  Maryam Amjad (Frontend Developer · Pakistan)\n• Ilham Kurnia Gustavakuan (Education Specialist · Indonesia)  |  Muhamad Rafi (Team Accountant)  |  Andini B. Soleman (Team Accountant)"
    p_t2.font.size = PPt(9.5)
    p_t2.font.color.rgb = PRGBColor(226, 232, 240)

    # SLIDE 2: The Problem
    s2 = create_base_slide("The Infodemic in an Age of AI: The Confidence-Ability Gap")
    
    def add_card(slide, left, top, width, height, title, body, tag=""):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = PRGBColor(30, 41, 59)
        card.line.color.rgb = PRGBColor(51, 65, 85)
        
        tx = slide.shapes.add_textbox(left + PInches(0.15), top + PInches(0.15), width - PInches(0.3), height - PInches(0.3))
        tf = tx.text_frame
        tf.word_wrap = True
        
        if tag:
            pt = tf.paragraphs[0]
            pt.text = tag.upper()
            pt.font.size = PPt(8.5)
            pt.font.bold = True
            pt.font.color.rgb = PRGBColor(13, 148, 136)
            p1 = tf.add_paragraph()
        else:
            p1 = tf.paragraphs[0]
            
        p1.text = title
        p1.font.size = PPt(13)
        p1.font.bold = True
        p1.font.color.rgb = PRGBColor(248, 250, 252)
        
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = PPt(10)
        p2.font.color.rgb = PRGBColor(203, 213, 225)
        return card

    add_card(s2, PInches(0.8), PInches(1.8), PInches(3.7), PInches(4.8), 
             "The Confidence-Ability Gap", 
             "62% of youth aged 8–17 feel confident navigating online spaces, yet lack practical cognitive skills to identify deepfakes, synthetic misinformation, or algorithmic ads (Ofcom, 2026).\n\nInformation fatigue pushes adolescents toward passive, short-form video & audio consumption.",
             "VULNERABILITY 01")
    
    add_card(s2, PInches(4.8), PInches(1.8), PInches(3.7), PInches(4.8), 
             "The AI 'Safety Vacuum'", 
             "Adolescents mistake simulated AI empathy for genuine understanding, becoming susceptible to ungrounded advice & synthetic bias (APA, 2026).\n\nStudents rely heavily on GenAI for schoolwork without receiving formal fact-checking or media literacy training (Rafique et al., 2026).",
             "VULNERABILITY 02")
    
    add_card(s2, PInches(8.8), PInches(1.8), PInches(3.7), PInches(4.8), 
             "The Text Literacy Barrier", 
             "700M+ low-literacy & rural citizens are excluded by text-heavy fact-checking sites in English.\n\nViral rumors spread through 30-second WhatsApp voice notes in local languages (Urdu, Indonesian, Spanish) where text tools completely fail.",
             "VULNERABILITY 03")

    # SLIDE 3: The 4 Core Objectives
    s3 = create_base_slide("VeriVoice Core Objectives: Transforming Youth into Active Guardians")
    
    objs = [
        ("Bridging the Digital Gap", "Create an accessible bridge between complex AI ecosystems and media literacy for youth consuming short-form feeds.", "GOAL 01"),
        ("Universal Oral Accessibility", "Provide an oral-first, voice-in/voice-out engine designed for low-literacy communities and text-fatigued users.", "GOAL 02"),
        ("Community Empowerment", "Transform users from passive consumers into active, critical agents of change and Community Information Stewards.", "GOAL 03"),
        ("Invisible Pedagogical Learning", "Embed Socratic explanations teaching WHY information is deceptive rather than delivering sterile True/False labels.", "GOAL 04")
    ]
    for idx, (otitle, obody, otag) in enumerate(objs):
        x = PInches(0.8 + (idx % 2) * 5.9)
        y = PInches(1.8 + (idx // 2) * 2.5)
        add_card(s3, x, y, PInches(5.7), PInches(2.3), otitle, obody, otag)

    # SLIDE 4: Architecture Pipeline
    s4 = create_base_slide("End-to-End Verification Pipeline: Bounded AI & Citation Integrity")
    
    steps = [
        ("1. Voice In", "Speechmatics & Groq Whisper ASR (Urdu, Indo, ES, EN)"),
        ("2. 15 Domains", "Health, Climate, AI Disinfo, Disaster, Space, Science, etc."),
        ("3. Authority", "WHO, NOAA, USGS, WMO, UNESCO Institutional Consensus"),
        ("4. Groq LPU", "Llama 3.3 70B in XML Isolation (<USER_CLAIM>, <EVIDENCE>)"),
        ("5. Guardrails", "CitationValidator Rejects Unretrieved URLs -> Safe UNCERTAIN"),
        ("6. Voice Out", "High-Definition Neural Audio (ElevenLabs & Edge-TTS)")
    ]
    for idx, (stitle, sdesc) in enumerate(steps):
        x = PInches(0.8 + (idx % 3) * 3.9)
        y = PInches(1.8 + (idx // 3) * 2.5)
        add_card(s4, x, y, PInches(3.7), PInches(2.3), stitle, sdesc, f"STAGE 0{idx+1}")

    # SLIDE 5: Socratic Dialogue & Innovation
    s5 = create_base_slide("Breakthrough Innovations: Socratic Dialogue & Invisible Learning")
    
    add_card(s5, PInches(0.8), PInches(1.8), PInches(5.7), PInches(4.8),
             "Socratic Dialogue vs Sterile Labels",
             "Traditional fact-checkers slap a clinical 'False' label that users ignore.\n\nVeriVoice converses like a critical mentor:\n“This audio claims the dam collapsed. However, official NDMA hydrological sensors record normal river flow. Notice how the forwarded voice message lacked a date, location, or source name—these are classic markers of urgency-engineered misinformation.”\n\nFosters internalized critical thinking without feeling like a classroom.",
             "INNOVATION 01")
    
    add_card(s5, PInches(6.8), PInches(1.8), PInches(5.7), PInches(4.8),
             "Cultural & Sensory Voice Adaptation",
             "Dismantles the 'technical literacy' barrier by translating dense WHO/NOAA/IPCC consensus into natural mother-tongue conversations (Urdu Nastaliq, Indonesian, Spanish, English).\n\nEmpowers rural mothers, students, and elders to verify rumors directly in community Discord channels and mobile web in under 1.8 seconds.",
             "INNOVATION 02")

    # SLIDE 6: Trust Nodes Network
    s6 = create_base_slide("Scalability: The Physical & Digital 'Trust Nodes Network'")
    
    add_card(s6, PInches(0.8), PInches(1.8), PInches(3.7), PInches(4.8),
             "Physical Trust Nodes",
             "Partnering with local pharmacies, libraries, schools, and community centers.\n\nPhysical placards provide instant QR-code 'Quick Start Kits' connecting directly to Web and messaging—zero app downloads needed.",
             "PILLAR 01")
    add_card(s6, PInches(4.8), PInches(1.8), PInches(3.7), PInches(4.8),
             "Hybrid Human-in-the-Loop",
             "When claims encounter extreme cultural nuance, queries escalate to a moderation dashboard managed by local youth leaders.\n\nCombines AI speed with authentic cultural context (Rafique et al., 2026).",
             "PILLAR 02")
    add_card(s6, PInches(8.8), PInches(1.8), PInches(3.7), PInches(4.8),
             "Frugal Economics",
             "Routine queries run on Groq LPU inference, costing < $0.0012 per query.\n\nScales seamlessly from one village to millions of users with multi-key API rotation.",
             "PILLAR 03")

    # SLIDE 7: Technical Proof
    s7 = create_base_slide("Empirical Feasibility: 22 Test Suites, 180 Tests & Sub-1.8s Latency")
    
    add_card(s7, PInches(0.8), PInches(1.8), PInches(3.7), PInches(4.8),
             "100% Green Test Suite",
             "• 22 Test Suites / 180 Automated Tests Passing (100% Green)\n• Specialized tests in tests/unescoAuthority.test.js & tests/chaosResilience.test.js\n• Full coverage of authority ranking, XML prompt isolation & URL allow-listing.",
             "SOFTWARE RIGOR")
    add_card(s7, PInches(4.8), PInches(1.8), PInches(3.7), PInches(4.8),
             "Sub-1.8s Response Latency",
             "• ASR Transcription: 380 ms\n• Domain Classifier: 2 ms\n• Institutional Retrieval: 450 ms\n• Bounded Groq LPU Reasoning: 420 ms\n• Neural TTS Synthesis: 350 ms\n• Total P95 Latency: ~1.60 – 1.85 seconds.",
             "PERFORMANCE BENCHMARK")
    add_card(s7, PInches(8.8), PInches(1.8), PInches(3.7), PInches(4.8),
             "Live Deployments",
             "• Production Frontend (Vercel): verivoice-ten.vercel.app\n• Cloud Backend API (Render): verivoice-unesco-hackathon.onrender.com\n• Discord Bot: Live 24/7 with 9 slash commands & audio processing.",
             "PRODUCTION READY")

    # SLIDE 8: SDGs & Sustainability
    s8 = create_base_slide("UN Sustainable Development Goals (SDGs) Alignment")
    
    add_card(s8, PInches(0.8), PInches(1.8), PInches(3.7), PInches(4.8),
             "SDG 4: Quality Education",
             "Expanding access to Media and Information Literacy (MIL) and AI literacy through inclusive voice-based learning experiences fostering critical thinking (UNESCO 2021, 2023).",
             "UN SDG 4")
    add_card(s8, PInches(4.8), PInches(1.8), PInches(3.7), PInches(4.8),
             "SDG 16: Peace & Justice",
             "Strengthening community resilience against hate speech, misinformation, and online panic while promoting social cohesion and informed dialogue (UN 2015).",
             "UN SDG 16")
    add_card(s8, PInches(8.8), PInches(1.8), PInches(3.7), PInches(4.8),
             "SDG 17: Partnerships",
             "Building a multi-stakeholder ecosystem connecting schools, youth organizations, pharmacies, libraries, and fact-checkers to scale media literacy globally.",
             "UN SDG 17")

    # SLIDE 9: Conclusion & Ask
    s9 = create_base_slide("The Resource Ask & Thessaloniki, Greece Vision")
    
    add_card(s9, PInches(0.8), PInches(1.8), PInches(5.7), PInches(4.8),
             "The Ask from UNESCO",
             "1. UNESCO MIL Alliance Partnership: Integration with UNESCO youth media hubs and global fact-checking observatories (EDMO, IFCN).\n\n2. Inference Grants: Subsidized API credits for Speechmatics, Groq, and WhatsApp Business API for 1,000,000 monthly free voice checks.\n\n3. Pilot Deployments: Facilitating physical Trust Node pilots in 10 community radio stations and rural hubs.",
             "RESOURCE ASK")
    add_card(s9, PInches(6.8), PInches(1.8), PInches(5.7), PInches(4.8),
             "Thessaloniki, Greece 2026",
             "“VeriVoice is an inclusive, voice-driven AI verification engine that empowers vulnerable youth and low-literacy communities to talk back to misinformation through Socratic dialogue and institutional truth in their native language.”\n\nJoin us at the Voice Festival in Thessaloniki, Greece, to bring the voice of truth to the next billion users!",
             "THE WINNING VISION")

    prs.save(output_path)
    print(f"[OK] Successfully generated PowerPoint Presentation: {output_path}")

if __name__ == '__main__':
    out_dir = os.path.join(os.getcwd(), 'docs', 'submission')
    os.makedirs(out_dir, exist_ok=True)
    
    docx_path = os.path.join(out_dir, 'VeriVoice_UNESCO_Written_Proposal_2026.docx')
    pptx_path = os.path.join(out_dir, 'VeriVoice_UNESCO_Pitch_Presentation_2026.pptx')
    
    build_word_document(docx_path)
    build_powerpoint_presentation(pptx_path)
